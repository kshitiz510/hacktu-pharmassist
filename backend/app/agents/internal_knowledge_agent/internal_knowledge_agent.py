from __future__ import annotations

import json
import io
from typing import Optional

from crewai import Agent, LLM

llm = LLM(model="groq/llama-3.3-70b-versatile", max_tokens=400)
llm_large = LLM(model="groq/llama-3.3-70b-versatile", max_tokens=4096)

internal_knowledge_agent = Agent(
    role="Internal Knowledge Agent",
    goal="Retrieve and synthesize internal documents, strategy decks, field insights, and research archives",
    backstory="""
    You are an expert internal knowledge analyst specializing in pharmaceutical strategic intelligence.
    
    Your responsibilities:
    - Retrieve and summarize internal documents (policy documents, strategy decks, technical notes)
    - Provide strategic synthesis and key insights
    - Identify important findings and recommendations
    
    OUTPUT FORMAT:
    - Start with "Based on internal policy documents and technical notes..."
    - Provide clear, well-structured descriptions
    - Highlight key findings and strategic implications
    """,
    tools=[],
    verbose=True,
    allow_delegation=False,
    llm=llm,
)

# Session-based document storage (in-memory, clears on server restart)
_session_documents: dict[str, dict] = {}


def store_document_for_session(session_id: str, filename: str, content: str, file_type: str) -> None:
    """Store parsed document content for a session."""
    _session_documents[session_id] = {
        "filename": filename,
        "content": content,
        "file_type": file_type,
    }
    print(f"[INTERNAL] Stored document '{filename}' for session {session_id} ({len(content)} chars)")


def get_document_for_session(session_id: str) -> Optional[dict]:
    """Retrieve stored document for a session."""
    return _session_documents.get(session_id)


def clear_document_for_session(session_id: str) -> None:
    """Clear document for a session."""
    if session_id in _session_documents:
        del _session_documents[session_id]


def parse_uploaded_file(file_content: bytes, filename: str) -> tuple[str, str]:
    """
    Parse uploaded file and extract text content.
    Supports: PDF, PPTX, XLSX, DOCX, TXT, CSV
    Returns: (extracted_text, file_type)
    """
    file_ext = filename.lower().split(".")[-1] if "." in filename else ""
    
    try:
        if file_ext == "pdf":
            return _parse_pdf(file_content), "pdf"
        elif file_ext == "pptx":
            return _parse_pptx(file_content), "pptx"
        elif file_ext in ["xlsx", "xls"]:
            return _parse_excel(file_content), "excel"
        elif file_ext == "docx":
            return _parse_docx(file_content), "docx"
        elif file_ext == "txt":
            return file_content.decode("utf-8", errors="ignore"), "txt"
        elif file_ext == "csv":
            return file_content.decode("utf-8", errors="ignore"), "csv"
        else:
            # Try to decode as text
            return file_content.decode("utf-8", errors="ignore"), "unknown"
    except Exception as e:
        print(f"[INTERNAL] Error parsing {filename}: {e}")
        return f"Error parsing file: {str(e)}", "error"


def _parse_pdf(content: bytes) -> str:
    """Parse PDF file and extract text."""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(io.BytesIO(content))
        text_parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        return "\n\n".join(text_parts)
    except Exception as e:
        return f"PDF parsing error: {str(e)}"


def _parse_pptx(content: bytes) -> str:
    """Parse PowerPoint file and extract text."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            text_parts.append("\n".join(slide_text))
        return "\n\n".join(text_parts)
    except Exception as e:
        return f"PPTX parsing error: {str(e)}"


def _parse_excel(content: bytes) -> str:
    """Parse Excel file and extract text."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(content), data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"--- Sheet: {sheet_name} ---"]
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                if any(row_text):
                    sheet_text.append(" | ".join(row_text))
            text_parts.append("\n".join(sheet_text))
        return "\n\n".join(text_parts)
    except Exception as e:
        return f"Excel parsing error: {str(e)}"


def _parse_docx(content: bytes) -> str:
    """Parse Word document and extract text."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        return "\n\n".join(text_parts)
    except Exception as e:
        return f"DOCX parsing error: {str(e)}"


def _llm_extract_params(user_prompt: str) -> tuple[str | None, str | None]:
    """Extract topic and focus area from user prompt."""
    prompt = f"""You are a pharmaceutical knowledge extraction expert. Extract the main topic and focus area from this query.

User Query: {user_prompt}

Return a JSON object with:
1. topic: The main pharmaceutical topic (drug name, therapy area, or general topic)
2. focus: The specific focus area (e.g., "regulatory", "clinical", "market", "strategic")

Return ONLY JSON, nothing else:"""

    try:
        raw = llm.call(messages=[{"role": "user", "content": prompt}])
    except Exception:
        try:
            raw = llm.call(prompt)
        except Exception:
            return None, None

    if not isinstance(raw, str):
        raw = str(raw)

    start = raw.find("{")
    end = raw.rfind("}")
    if start == -1 or end == -1:
        return None, None

    try:
        data = json.loads(raw[start:end + 1])
        topic = data.get("topic")
        focus = data.get("focus")
        return topic, focus
    except Exception:
        return None, None


def _llm_analyze_document(document_content: str, user_query: str, filename: str) -> dict:
    """Use LLM to analyze uploaded document content."""
    # Truncate document if too long
    max_content_length = 12000
    truncated = document_content[:max_content_length]
    if len(document_content) > max_content_length:
        truncated += f"\n\n[Document truncated - showing first {max_content_length} characters of {len(document_content)} total]"

    prompt = f"""You are a pharmaceutical internal knowledge analyst. Analyze this internal document and provide insights relevant to the user's query.

DOCUMENT NAME: {filename}
USER QUERY: {user_query}

DOCUMENT CONTENT:
{truncated}

Provide a comprehensive analysis with this EXACT structure:
1. Start your response with "Based on internal policy documents and technical notes..."
2. Document Overview: Brief summary of what the document contains
3. Key Findings: 3-5 most important findings relevant to the query
4. Strategic Implications: How these findings impact pharmaceutical strategy
5. Recommendations: 2-3 actionable recommendations based on the document
6. Important Data Points: Any specific numbers, dates, or metrics mentioned

Return as JSON:
{{
    "overview": "Based on internal policy documents and technical notes, [summary]",
    "key_findings": ["finding 1", "finding 2", ...],
    "strategic_implications": "description of strategic implications",
    "recommendations": ["rec 1", "rec 2", ...],
    "data_points": ["data point 1", "data point 2", ...],
    "confidence": "high/medium/low"
}}"""

    try:
        raw = llm_large.call(messages=[{"role": "user", "content": prompt}])
    except Exception:
        try:
            raw = llm_large.call(prompt)
        except Exception as e:
            print(f"[INTERNAL] LLM analysis failed: {e}")
            return None

    if not isinstance(raw, str):
        raw = str(raw)

    start = raw.find("{")
    end = raw.rfind("}")
    if start == -1 or end == -1:
        return None

    try:
        return json.loads(raw[start:end + 1])
    except Exception:
        return None


def _llm_generate_internal_knowledge(user_query: str, topic: str | None) -> dict:
    """Generate internal knowledge insights using LLM when no document is provided."""
    prompt = f"""You are a pharmaceutical internal knowledge analyst with access to internal databases, policy documents, and technical notes.

USER QUERY: {user_query}
TOPIC: {topic or "General Pharmaceutical Intelligence"}

Generate comprehensive internal knowledge insights as if retrieved from internal company databases. Start with "Based on internal policy documents and technical notes..."

Provide insights in this EXACT JSON format:
{{
    "overview": "Based on internal policy documents and technical notes, [comprehensive overview of the topic from internal knowledge perspective]",
    "key_findings": [
        "Internal finding 1 with specific details",
        "Internal finding 2 with strategic context",
        "Internal finding 3 with market implications",
        "Internal finding 4 with regulatory considerations",
        "Internal finding 5 with competitive intelligence"
    ],
    "strategic_implications": "Detailed analysis of how these internal insights impact pharmaceutical strategy, including market positioning, R&D priorities, and competitive dynamics",
    "recommendations": [
        "Strategic recommendation 1 based on internal analysis",
        "Tactical recommendation 2 for immediate action",
        "Long-term recommendation 3 for portfolio planning"
    ],
    "internal_references": [
        "Internal Policy Document: Pharmaceutical Development Guidelines v2.3",
        "Technical Note: Market Intelligence Report Q4 2025",
        "Strategy Deck: Competitive Landscape Analysis"
    ],
    "confidence": "medium",
    "data_source": "Internal Knowledge Database"
}}

Make the insights realistic, detailed, and professionally written as if from actual internal pharmaceutical company documents."""

    try:
        raw = llm_large.call(messages=[{"role": "user", "content": prompt}])
    except Exception:
        try:
            raw = llm_large.call(prompt)
        except Exception as e:
            print(f"[INTERNAL] LLM generation failed: {e}")
            return None

    if not isinstance(raw, str):
        raw = str(raw)

    start = raw.find("{")
    end = raw.rfind("}")
    if start == -1 or end == -1:
        return None

    try:
        return json.loads(raw[start:end + 1])
    except Exception:
        return None


def run_internal_knowledge_agent(
    user_prompt: str,
    session_id: str | None = None,
) -> dict:
    """
    Run the Internal Knowledge Agent.
    - If a document is uploaded for the session, analyze it
    - Otherwise, use LLM to generate internal knowledge insights
    """
    print(f"[INTERNAL] Starting agent with prompt: {user_prompt}")
    print(f"[INTERNAL] Session ID: {session_id}")

    # Extract topic from query
    topic, focus = _llm_extract_params(user_prompt)
    print(f"[INTERNAL] Extracted topic: {topic}, focus: {focus}")

    # Check for uploaded document
    document = get_document_for_session(session_id) if session_id else None

    if document:
        # Analyze the uploaded document
        print(f"[INTERNAL] Analyzing uploaded document: {document['filename']}")
        
        analysis = _llm_analyze_document(
            document["content"],
            user_prompt,
            document["filename"]
        )

        if analysis:
            payload = {
                "source": "uploaded_document",
                "filename": document["filename"],
                "file_type": document["file_type"],
                "analysis": analysis,
                "overview": analysis.get("overview", "Based on internal policy documents and technical notes, analysis completed."),
                "key_findings": analysis.get("key_findings", []),
                "strategic_implications": analysis.get("strategic_implications", ""),
                "recommendations": analysis.get("recommendations", []),
                "data_points": analysis.get("data_points", []),
                "confidence": analysis.get("confidence", "medium"),
            }
            
            # Build canonical summary for banner
            key_findings_count = len(analysis.get("key_findings", []))
            has_recommendations = len(analysis.get("recommendations", [])) > 0
            confidence_level = analysis.get("confidence", "medium")
            
            banner_summary = {
                "researcherQuestion": "Does internal knowledge support this research?",
                "answer": "Yes" if key_findings_count > 0 else "Limited",
                "explainers": []
            }
            
            if key_findings_count > 0:
                banner_summary["explainers"].append(f"{key_findings_count} key findings")
            if has_recommendations:
                banner_summary["explainers"].append("Strategic recommendations available")
            if filename:
                banner_summary["explainers"].append(f"Source: {filename}")
            
            payload["summary"] = banner_summary
            
            # Generate suggested next prompts
            suggested_next_prompts = [
                {"prompt": f"Show clinical trials for {topic or 'this drug'}"},
                {"prompt": f"Analyze patent landscape for {topic or 'this compound'}"},
                {"prompt": f"Get market intelligence for {topic or 'this therapy area'}"}
            ]
            payload["suggestedNextPrompts"] = suggested_next_prompts

            return {
                "status": "success",
                "data": payload,
                "visualizations": _build_internal_knowledge_visualizations(payload),
            }
    
    # No document - use LLM fallback
    print(f"[INTERNAL] No document uploaded, using LLM to generate internal knowledge")
    
    llm_data = _llm_generate_internal_knowledge(user_prompt, topic)

    if llm_data:
        # Build canonical summary for banner
        key_findings_count = len(llm_data.get("key_findings", []))
        has_recommendations = len(llm_data.get("recommendations", [])) > 0
        
        banner_summary = {
            "researcherQuestion": "Does internal knowledge support this research?",
            "answer": "Yes" if key_findings_count > 0 else "Limited",
            "explainers": []
        }
        
        if key_findings_count > 0:
            banner_summary["explainers"].append(f"{key_findings_count} key findings")
        if has_recommendations:
            banner_summary["explainers"].append("Strategic recommendations available")
        banner_summary["explainers"].append("Source: Internal Knowledge Database")
        
        # Generate suggested next prompts
        suggested_next_prompts = [
            {"prompt": f"Show clinical trials for {topic or 'this drug'}"},
            {"prompt": f"Analyze patent landscape for {topic or 'this compound'}"},
            {"prompt": f"Get market intelligence for {topic or 'this therapy area'}"}
        ]
        
        payload = {
            "source": "internal_database",
            "filename": None,
            "file_type": None,
            "analysis": llm_data,
            "overview": llm_data.get("overview", "Based on internal policy documents and technical notes, insights retrieved from internal knowledge base."),
            "key_findings": llm_data.get("key_findings", []),
            "strategic_implications": llm_data.get("strategic_implications", ""),
            "recommendations": llm_data.get("recommendations", []),
            "internal_references": llm_data.get("internal_references", []),
            "confidence": llm_data.get("confidence", "medium"),
            "data_source": llm_data.get("data_source", "Internal Knowledge Database"),
            "summary": banner_summary,
            "suggestedNextPrompts": suggested_next_prompts,
        }

        return {
            "status": "success",
            "data": payload,
            "visualizations": _build_internal_knowledge_visualizations(payload),
        }

    return {
        "status": "error",
        "message": "Failed to retrieve internal knowledge insights",
    }


def _build_internal_knowledge_visualizations(payload: dict) -> list:
    """Build visualization payload for Internal Knowledge Agent with enhanced formatting."""
    viz = []
    
    source = payload.get("source", "internal_database")
    filename = payload.get("filename")
    confidence = payload.get("confidence", "medium")
    
    # Title based on source
    if source == "uploaded_document" and filename:
        title = f"📄 Document Analysis: {filename}"
    else:
        title = "📚 Internal Knowledge Insights"
    
    # NOTE: Top summary card removed per design spec (compact template)
    
    # 1. Document Title (one line) - No card, just metadata
    # Title already set above as: f"📄 Document Analysis: {filename}" or "📚 Internal Knowledge Insights"
    
    # 2. Executive Summary (one-line sentence) - Clean text, no markdown
    overview = payload.get("overview", "")
    executive_summary = overview.split(".")[0] + "." if overview else "No summary available"
    
    # 3. Key Findings - Plain text (no ** markdown - causes display issues)
    key_findings = payload.get("key_findings", [])
    if key_findings:
        # Limit to 3-6 findings, NO word truncation - show full text
        top_findings = key_findings[:6]
        
        # Plain text format - each finding on its own line
        findings_lines = []
        for i, finding in enumerate(top_findings):
            findings_lines.append(f"{i+1}. {finding}")
        findings_text = "\n\n".join(findings_lines)  # Double newline for better spacing
        
        viz.append({
            "id": "key_findings",
            "vizType": "text",
            "title": title,
            "description": findings_text,
            "data": {
                "content": findings_text, 
                "items": top_findings,
                "format": "plain",
                "executive_summary": executive_summary
            },
        })
    
    # 3. Strategic Context (2-3 paragraphs) - Before Quick Facts
    implications = payload.get("strategic_implications", "")
    if implications:
        # Split into paragraphs, keep up to 3
        paragraphs = implications.split("\n\n")
        short_paragraphs = paragraphs[:3]  # No word truncation
        
        strategic_context = "\n\n".join(short_paragraphs)
        
        viz.append({
            "id": "strategic_context",
            "vizType": "text",
            "title": "💡 Strategic Context",
            "description": strategic_context,
            "data": {
                "content": strategic_context,
                "format": "paragraph"
            },
        })
    
    # 4. Quick Facts (metrics row: #findings | Confidence | Downloads) - Inline format
    data_points = payload.get("data_points", [])
    quick_facts_inline = f"{len(key_findings)} findings | Confidence: {confidence.title()}"
    if filename:
        quick_facts_inline += f" | Source: {filename}"
    
    viz.append({
        "id": "quick_facts",
        "vizType": "metric",
        "title": "📊 Quick Facts",
        "description": quick_facts_inline,
        "data": {
            "findings_count": len(key_findings),
            "confidence": confidence,
            "source": filename if filename else "Internal Database",
            "format": "inline"
        },
    })
    
    # 5. Recommended Actions (plain text, no markdown)
    recommendations = payload.get("recommendations", [])
    if recommendations:
        # Simple numbered list - no "Immediate/Near-term/Long-term" prefixes
        rec_lines = []
        for i, rec in enumerate(recommendations[:3]):
            rec_lines.append(f"{i+1}. {rec}")
        
        rec_text = "\n".join(rec_lines)
        
        viz.append({
            "id": "recommended_actions",
            "vizType": "text",
            "title": "Recommended Actions",
            "description": rec_text,
            "data": {
                "content": rec_text, 
                "items": recommendations[:3],
                "format": "plain"
            },
        })
    
    # 6. Evidence & Sources (sources stored in data, NOT in description to avoid duplication)
    internal_refs = payload.get("internal_references", [])
    sources = []
    if filename:
        sources.append(filename)
    if internal_refs:
        for ref in internal_refs[:5]:
            if isinstance(ref, dict):
                source_name = ref.get("title", "Unknown")
                # Skip if source_name already in sources (dedup)
                if source_name not in sources:
                    sources.append(source_name)
            else:
                ref_str = str(ref)
                if ref_str not in sources:
                    sources.append(ref_str)
    
    if sources:
        # Only show sources list once (in data.sources), not in description
        viz.append({
            "id": "evidence_sources",
            "vizType": "sources",  # Changed vizType for proper rendering
            "title": "Evidence & Sources",
            "description": "",  # Empty to avoid duplication
            "data": {
                "sources": sources,
                "format": "list"
            },
        })
    
    # 7. UI Affordances - Download CSV and Expand View buttons
    viz.append({
        "id": "ui_actions",
        "vizType": "actions",
        "title": "",
        "description": "",
        "data": {
            "actions": [
                {"label": "Download CSV", "action": "download_csv", "enabled": True},
                {"label": "Expand View", "action": "expand_view", "enabled": True}
            ]
        },
    })
    
    # 8. Data Points (if from document, keep brief - plain text, no markdown)
    if data_points:
        dp_text = " | ".join(data_points[:5])
        if len(data_points) > 5:
            dp_text += f" | ...and {len(data_points) - 5} more"
        
        viz.append({
            "id": "data_points",
            "vizType": "text",
            "title": "Key Metrics & Data",
            "description": dp_text,
            "data": {
                "content": dp_text, 
                "items": data_points,
                "format": "plain"
            },
        })
    
    return viz
